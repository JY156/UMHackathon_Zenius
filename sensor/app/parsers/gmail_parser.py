# sensor/app/parsers/gmail_parser.py
import base64
import io
from bs4 import BeautifulSoup

# ✅ Import optional libraries (install only what you need)
# pip install pypdf python-docx python-pptx openpyxl pytesseract pillow

# PDF support
try:
    from pypdf import PdfReader
    PDF_LIB = "pypdf"
except ImportError:
    try:
        import PyPDF2
        PDF_LIB = "PyPDF2"
    except ImportError:
        PDF_LIB = None

# DOCX (Word) support
try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

# PPTX (PowerPoint) support
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# XLSX (Excel) support
try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

# PNG/JPG OCR support (OPTIONAL - requires Tesseract OCR installed on system)
try:
    import pytesseract
    from PIL import Image
    HAS_OCR = True
except ImportError:
    HAS_OCR = False
    print("⚠️ OCR not available - install pytesseract and pillow for image text extraction")


def clean_html(raw_html):
    if not raw_html:
        return ""
    soup = BeautifulSoup(raw_html, "html.parser")
    for script_or_style in soup(["script", "style"]):
        script_or_style.decompose()
    return soup.get_text(separator=' ').strip()


def extract_attachment_text(filename: str, base64_data: str, mime_type: str) -> str:
    """
    ✅ Universal extractor: handles PDF, DOCX, PPTX, XLSX, TXT, and images (OCR).
    Returns extracted text or placeholder if unsupported/failed.
    """
    ext = filename.lower().split('.')[-1] if '.' in filename else ""
    
    # Decode base64url to bytes (Gmail uses base64url encoding)
    try:
        normalized = base64_data.replace('-', '+').replace('_', '/')
        padding = len(normalized) % 4
        if padding:
            normalized += '=' * (4 - padding)
        file_bytes = base64.b64decode(normalized)
    except Exception as e:
        return f"[Attachment decode error: {filename}]"
    
    try:
        # ✅ PDF support
        if ext == 'pdf' or mime_type == 'application/pdf':
            if not PDF_LIB:
                return "[PDF: Install 'pip install pypdf']"
            
            if PDF_LIB == "pypdf":
                reader = PdfReader(io.BytesIO(file_bytes))
            else:
                reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
            
            text_parts = []
            for page in reader.pages[:5]:  # First 5 pages
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    text_parts.append(f"[Page {reader.pages.index(page)+1}] {page_text.strip()}")
            
            extracted = "\n\n".join(text_parts)
            return extracted[:3000] + "..." if len(extracted) > 3000 else extracted
        
        # ✅ DOCX (Word) support
        elif ext in ['doc', 'docx'] or 'word' in mime_type or 'document' in mime_type:
            if not HAS_DOCX:
                return "[DOCX: Install 'pip install python-docx']"
            
            doc = Document(io.BytesIO(file_bytes))
            paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            
            # Extract from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            paragraphs.append(f"[Table] {cell.text.strip()}")
            
            extracted = "\n\n".join(paragraphs)
            return extracted[:3000] + "..." if len(extracted) > 3000 else extracted
        
        # ✅ PPTX (PowerPoint) support
        elif ext in ['ppt', 'pptx'] or 'presentation' in mime_type or 'powerpoint' in mime_type:
            if not HAS_PPTX:
                return "[PPTX: Install 'pip install python-pptx']"
            
            prs = Presentation(io.BytesIO(file_bytes))
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    # Extract text from tables in slides
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    slide_text.append(f"[Table] {cell.text.strip()}")
                
                if slide_text:
                    text_parts.append(f"[Slide {slide_num+1}] " + " | ".join(slide_text))
            
            extracted = "\n\n".join(text_parts)
            return extracted[:3000] + "..." if len(extracted) > 3000 else extracted
        
        # ✅ XLSX (Excel) support
        elif ext in ['xls', 'xlsx', 'csv'] or 'excel' in mime_type or 'spreadsheet' in mime_type:
            if ext == 'csv':
                return file_bytes.decode('utf-8', errors='ignore')[:3000]
            
            if not HAS_XLSX:
                return "[XLSX: Install 'pip install openpyxl']"
            
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
            text_parts = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    # Filter out None/empty cells
                    clean_row = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                    if clean_row:
                        rows.append(" | ".join(clean_row))
                if rows:
                    text_parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows[:15]))  # First 15 rows
            
            extracted = "\n\n".join(text_parts)
            return extracted[:3000] + "..." if len(extracted) > 3000 else extracted
        
        # ✅ Plain text files (TXT, MD, LOG, RTF, CSV)
        elif ext in ['txt', 'md', 'log', 'rtf', 'csv']:
            return file_bytes.decode('utf-8', errors='ignore')[:3000]
        
        # ⚠️ PNG/JPG OCR support (OPTIONAL)
        elif ext in ['png', 'jpg', 'jpeg', 'gif', 'bmp'] or 'image' in mime_type:
            if not HAS_OCR:
                return f"[Image: {filename} - OCR not available. Install pytesseract + Tesseract OCR]"
            
            try:
                img = Image.open(io.BytesIO(file_bytes))
                extracted_text = pytesseract.image_to_string(img)
                return extracted_text.strip()[:2000] if extracted_text.strip() else "[Image: No text detected]"
            except Exception as e:
                return f"[Image OCR failed: {str(e)}]"
        
        # ❌ Unsupported type - return placeholder
        else:
            return f"[Attachment: {filename} (.{ext}) - text extraction not supported]"
            
    except Exception as e:
        print(f"⚠️ Extraction failed for {filename}: {type(e).__name__}: {str(e)}")
        return f"[Attachment: {filename} - extraction failed: {type(e).__name__}]"


def parse_gmail_message(service, message_raw: dict):
    payload = message_raw.get('payload', {})
    headers = payload.get('headers', [])
    message_id = message_raw.get('id')

    def get_header_value(name):
        return next((h['value'] for h in headers if h['name'].lower() == name.lower()), None)
    
    sender = get_header_value('From') or "Unknown"
    subject = get_header_value('Subject') or "No Subject"
    thread_id = message_raw.get('threadId')

    parts = payload.get('parts', [])
    body = ""
    attachments = []

    def walk_parts(all_parts):
        nonlocal body
        for part in all_parts:
            filename = part.get('filename')
            mime_type = part.get('mimeType')
            
            if filename:
                attach_id = part.get('body', {}).get('attachmentId')
                if attach_id:
                    attachment_meta = service.users().messages().attachments().get(
                        userId='me', messageId=message_id, id=attach_id
                    ).execute()
                    
                    base64_data = attachment_meta.get('data', '')
                    
                    # ✅ USE THE UNIVERSAL EXTRACTOR
                    extracted_text = extract_attachment_text(filename, base64_data, mime_type)
                    
                    # Log extraction result
                    if extracted_text and not extracted_text.startswith("["):
                        print(f"📎 ✅ {filename}: {len(extracted_text)} chars extracted")
                    else:
                        print(f"📎 ⚠️ {filename}: {extracted_text}")
                    
                    attachments.append({
                        "filename": filename,
                        "mimeType": mime_type,
                        "data": base64_data,
                        "size": part.get('body', {}).get('size'),
                        "attachmentId": attach_id,
                        "extracted_text": extracted_text,
                        "extraction_success": not extracted_text.startswith("[")
                    })
            
            elif mime_type in ['text/plain', 'text/html'] and not body:
                data = part.get('body', {}).get('data', "")
                if data:
                    body = data
            
            if 'parts' in part:
                walk_parts(part['parts'])

    if not parts:
        body = payload.get('body', {}).get('data', "")
    else:
        walk_parts(parts)

    decoded_body = base64.urlsafe_b64decode(body).decode('utf-8') if body else ""
    cleaned_body = clean_html(decoded_body)

    return {
        "thread_id": thread_id,
        "subject": subject,
        "raw_body": decoded_body,
        "cleaned_body": cleaned_body,
        "user_id": sender,
        "timestamp": message_raw.get("internalDate"),
        "source": "gmail",
        "attachments": attachments
    }